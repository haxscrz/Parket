from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
PRIMARY_COLOR = RGBColor(0, 102, 204)  # Blue
SECONDARY_COLOR = RGBColor(255, 153, 0)  # Orange
DARK_TEXT = RGBColor(51, 51, 51)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(title, subtitle, logo_path=None):
    """Add a custom title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Add logo if available
    if logo_path and os.path.exists(logo_path):
        left = Inches(3.5)
        top = Inches(1.5)
        slide.shapes.add_picture(logo_path, left, top, width=Inches(3))
        title_top = Inches(3.8)
    else:
        title_top = Inches(2.5)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), title_top, Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(200, 220, 255)
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, points):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Add underline
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(5), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY_COLOR
    shape.line.color.rgb = SECONDARY_COLOR
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_TEXT
        p.level = 0
        p.space_before = Pt(12)
        p.space_after = Pt(12)

def add_two_column_slide(title, left_title, left_points, right_title, right_points):
    """Add a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # Add underline
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(5), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY_COLOR
    shape.line.color.rgb = SECONDARY_COLOR
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    p = left_frame.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    for point in left_points:
        p = left_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    p = right_frame.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    for point in right_points:
        p = right_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(8)

# Slide 1: Title Slide
logo_path = r"d:\DOWNLOADS\Parket\payment_logo\gcash.png"
add_title_slide("🅿️ PARKET", "AI-Powered Smart Parking Solution")

# Slide 2: What is Parket
add_content_slide("What is Parket?", [
    "🎯 An intelligent parking app powered by Artificial Intelligence and Machine Learning",
    "📊 Helps drivers find optimal parking spots with real-time predictions",
    "🔍 Analyzes historical occupancy data to forecast parking availability",
    "⚡ Reduces time wasted searching for parking spots",
    "🌍 Works seamlessly across shopping malls and establishments"
])

# Slide 3: Problem Statement
add_content_slide("The Problem We Solve", [
    "🚗 Drivers waste 30% of driving time searching for parking spots",
    "😤 Unpredictable parking availability causes frustration",
    "⏰ Peak hours lead to congestion and longer wait times",
    "💰 Unnecessary fuel consumption and increased emissions",
    "📍 No data-driven insights for parking decisions"
])

# Slide 4: Features
add_two_column_slide("Key Features",
    "AI & Machine Learning",
    [
        "• Peak Hours Prediction",
        "• Best Time Suggestions",
        "• Smart Slot Routing",
        "• Occupancy Forecasts"
    ],
    "User Features",
    [
        "• Real-time Availability",
        "• E-wallet Integration",
        "• Parking History",
        "• Dark/Light Mode"
    ]
)

# Slide 5: How It Was Made
add_two_column_slide("Technology Stack",
    "Frontend",
    [
        "• React 18 + TypeScript",
        "• Vite (Fast Builds)",
        "• Tailwind CSS",
        "• Framer Motion"
    ],
    "Backend & ML",
    [
        "• TensorFlow.js",
        "• Python (Training)",
        "• Random Forest (94.5% Accuracy)",
        "• Real-time Data API"
    ]
)

# Slide 6: Real-Life Use Cases
add_content_slide("Real-Life Applications", [
    "🛍️ Shopping Malls: Guide shoppers to available spots instantly",
    "🏢 Office Buildings: Predict parking availability during peak hours",
    "🅿️ Parking Lots: Optimize occupancy and revenue",
    "🌆 Urban Areas: Reduce traffic congestion and emissions",
    "🔄 Integration: Works with existing parking management systems"
])

# Slide 7: Revenue Model
add_content_slide("Revenue Model - Subscription Service", [
    "💵 Establishments & Malls pay monthly subscription fees",
    "📈 Pricing tiers based on parking lot size and user traffic",
    "🎯 Malls benefit from: Better customer experience & increased visits",
    "💡 Increased customer satisfaction = Higher retention & spending",
    "📊 Data insights for mall operators to optimize operations"
])

# Slide 8: Demo Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
title_frame = title_box.text_frame
p = title_frame.paragraphs[0]
p.text = "🎬 APP DEMO"
p.font.size = Pt(88)
p.font.bold = True
p.font.color.rgb = PRIMARY_COLOR
p.alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
p = subtitle_frame.paragraphs[0]
p.text = "Live demonstration of the Parket app interface, features, and AI predictions"
p.font.size = Pt(22)
p.font.color.rgb = SECONDARY_COLOR
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"d:\DOWNLOADS\Parket\Parket_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation created successfully: {output_path}")
